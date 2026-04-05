import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- THEME COLORS ---
BG_COLOR = RGBColor(12, 16, 28)           # Deep Navy Dark
BOX_BG = RGBColor(25, 30, 48)             # Panel background
BOX_BORDER = RGBColor(60, 75, 110)        # Panel border
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GREY = RGBColor(180, 190, 210)
GOLD = RGBColor(255, 215, 0)
GREEN = RGBColor(0, 200, 100)
RED = RGBColor(220, 50, 50)
BLUE = RGBColor(0, 150, 255)

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def create_badge(slide, text, left, top, width, height, bg_color, text_color):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.width = 0
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return badge

def create_info_panel(slide, emoji, title, value, subtitle, left, top):
    width = Inches(2.2)
    height = Inches(1.5)
    
    # Background Panel
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = BOX_BG
    box.line.color.rgb = BOX_BORDER
    box.line.width = Pt(1.5)

    # Emoji Box
    e_box = slide.shapes.add_textbox(left, top + Inches(0.05), width, Inches(0.4))
    e_p = e_box.text_frame.paragraphs[0]
    e_p.text = emoji
    e_p.font.size = Pt(18)
    e_p.alignment = PP_ALIGN.CENTER

    # Title
    t_box = slide.shapes.add_textbox(left, top + Inches(0.35), width, Inches(0.3))
    t_p = t_box.text_frame.paragraphs[0]
    t_p.text = title
    t_p.font.name = 'Arial'
    t_p.font.size = Pt(10)
    t_p.font.bold = True
    t_p.font.color.rgb = TEXT_GREY
    t_p.alignment = PP_ALIGN.CENTER

    # Value
    v_box = slide.shapes.add_textbox(left, top + Inches(0.6), width, Inches(0.5))
    v_p = v_box.text_frame.paragraphs[0]
    v_p.text = str(value)
    v_p.font.name = 'Arial Black'
    v_p.font.size = Pt(20)
    v_p.font.color.rgb = GOLD
    v_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    s_box = slide.shapes.add_textbox(left, top + Inches(1.1), width, Inches(0.3))
    s_p = s_box.text_frame.paragraphs[0]
    s_p.text = subtitle
    s_p.font.name = 'Arial'
    s_p.font.size = Pt(8)
    s_p.font.color.rgb = TEXT_GREY
    s_p.alignment = PP_ALIGN.CENTER

def generate_ppt():
    players_json_path = 'client/public/players.json'
    images_dir = 'client/public/players'
    output_path = 'IPL_Final_Match_Theme_v6.pptx'

    with open(players_json_path, 'r', encoding='utf-8') as f:
        players = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Calculate total for progress bar
    total_players = len(players)

    for idx, player in enumerate(players):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, BG_COLOR)

        # --- HEADER ROW ---
        h_tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(4), Inches(0.4)).text_frame
        hp = h_tf.paragraphs[0]
        hp.text = "IPL 2026 COLLEGE AUCTION"
        hp.font.size = Pt(14)
        hp.font.bold = True
        hp.font.color.rgb = TEXT_WHITE

        # Top Right Progress
        pct = int(((idx + 1) / total_players) * 100)
        p_tf = slide.shapes.add_textbox(Inches(10), Inches(0.3), Inches(3), Inches(0.4)).text_frame
        pp = p_tf.paragraphs[0]
        pp.text = f"PLAYER #{idx+1} OF {total_players}  |  {pct}% Complete"
        pp.font.size = Pt(11)
        pp.font.color.rgb = TEXT_GREY
        pp.alignment = PP_ALIGN.RIGHT

        # --- TOP BADGES ---
        left_pos = 0.5
        role = player.get('role', 'All-Rounder').upper()
        role_emoji = "🏏" if "BAT" in role else ("🥎" if "BOWL" in role else ("🧤" if "WICKET" in role else "⚔️"))
        create_badge(slide, f"{role_emoji}  {role}", Inches(left_pos), Inches(1), Inches(1.8), Inches(0.3), BLUE, TEXT_WHITE)
        left_pos += 2.0

        country = player.get('country', 'INDIA').upper()
        create_badge(slide, f"🇮🇳  {country}", Inches(left_pos), Inches(1), Inches(1.3), Inches(0.3), BOX_BG, TEXT_WHITE)
        left_pos += 1.5

        create_badge(slide, "✓  AVAILABLE", Inches(left_pos), Inches(1), Inches(1.5), Inches(0.3), GREEN, TEXT_WHITE)
        
        # --- PLAYER NAME ---
        name_tf = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(8), Inches(1.2)).text_frame
        np = name_tf.paragraphs[0]
        np.text = player['name'].upper()
        np.font.name = 'Arial Black'
        np.font.size = Pt(56)
        np.font.color.rgb = TEXT_WHITE

        # --- 3 DATA PANELS ---
        # 1. Base Price
        create_info_panel(slide, "💰", "BASE PRICE", player.get('basePrice', 'N/A'), "Auction Starting Bid", Inches(0.5), Inches(3.0))
        
        # 2. Rating
        rating_val = int(player.get('rating', 0))
        create_info_panel(slide, "⭐", "RATING", f"{rating_val}/10", "Player Skill Rating", Inches(2.9), Inches(3.0))

        # 3. Origin (Centered gracefully underneath the top two panels)
        create_info_panel(slide, "🌍", "ORIGIN", country, "Home Country", Inches(1.7), Inches(4.7))

        # --- STAR RATING GRAPHIC BELOW PANELS ---
        stars = "★" * rating_val + "☆" * (10 - rating_val)
        st_tf = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(5), Inches(0.5)).text_frame
        stp = st_tf.paragraphs[0]
        stp.text = f"Skill Level: {stars}"
        stp.font.size = Pt(20)
        stp.font.color.rgb = GOLD

        # --- PLAYER IMAGE (RIGHT SIDE) ---
        img_name = player.get('image_file') or f"{player['id']}.png"
        img_path = os.path.join(images_dir, img_name)
        if os.path.exists(img_path):
            try:
                # Add an elegant glowing circle behind the player
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(1.0), Inches(6), Inches(6))
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(18, 25, 45) # Slightly lighter deep blue
                circle.line.width = 0
                
                pic = slide.shapes.add_picture(img_path, Inches(6.8), Inches(1.2), height=Inches(5.8))
            except: pass

    prs.save(output_path)
    print(f"Final Matched Theme Generated: {output_path}")

if __name__ == "__main__":
    generate_ppt()
