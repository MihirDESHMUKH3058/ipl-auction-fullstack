import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Constants & Theme Colors ---
NAVY_BLUE = RGBColor(10, 17, 40)      # #0a1128
GOLDEN_YELLOW = RGBColor(253, 185, 19) # #fdb913
WHITE = RGBColor(255, 255, 255)       # #ffffff
DARK_GREY_NAVY = RGBColor(30, 41, 59) # #1e293b
BORDER_COLOR = RGBColor(51, 65, 85)     # #334155

def extract_mock_stats(file_path):
    if not os.path.exists(file_path): return {}
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    match = re.search(r'export const mockPlayerStats = ({[\s\S]*?});', content)
    if not match: return {}
    js_obj = match.group(1)
    json_str = re.sub(r'(\s)(\w+):', r'\1"\2":', js_obj)
    json_str = re.sub(r',\s*}', '}', json_str)
    try: return json.loads(json_str)
    except: return {}

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_rules_box(slide, left, top, width, height, title, items):
    # Rule box background
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_GREY_NAVY
    rect.line.color.rgb = BORDER_COLOR
    rect.line.width = Pt(1)

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLDEN_YELLOW

    # Items
    content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for item in items:
        p = ctf.add_paragraph()
        p.text = f"▶  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

def generate_ppt():
    players_json_path = 'client/public/players.json'
    mock_stats_path = 'client/src/mockStats.js'
    images_dir = 'client/public/players'
    output_path = 'IPL_2026_Auction_Catalog.pptx'

    with open(players_json_path, 'r', encoding='utf-8') as f:
        players = json.load(f)
    stats_data = extract_mock_stats(mock_stats_path)

    ind_count = len([p for p in players if p.get('overseas') == 'Indian'])
    ovs_count = len([p for p in players if p.get('overseas') == 'Overseas'])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Professional Intro ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, NAVY_BLUE)

    # Left Sidebar Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0), Inches(0.05), Inches(7.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLDEN_YELLOW
    line.line.width = 0

    # Left Sidebar Content
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(3.5), Inches(1))
    tf = logo_box.text_frame
    tf.text = "IPL 2026"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = GOLDEN_YELLOW
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(3.5), Inches(1))
    tf = sub_box.text_frame
    tf.text = "COLLEGE AUCTION\nOFFICIAL PLAYER CATALOG"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Right Main Content
    cat_lbl = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(8), Inches(0.5))
    tf = cat_lbl.text_frame
    tf.text = "P L A Y E R   C A T A L O G"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = GOLDEN_YELLOW

    main_title = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(8), Inches(2))
    tf = main_title.text_frame
    tf.text = "IPL 2026 College\nAuction"
    p = tf.paragraphs[0]
    p.font.size = Pt(88)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # KPI Boxes
    kpis = [
        (f"{len(players)}", "Total Players"),
        (f"{ind_count}", "Indian Players"),
        (f"{ovs_count}", "Overseas Players"),
        ("10", "IPL Teams")
    ]
    for i, (val, lbl) in enumerate(kpis):
        x = Inches(5 + (i * 2.1))
        y = Inches(3.8)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.9), Inches(1.6))
        rect.line.color.rgb = GOLDEN_YELLOW
        rect.fill.background()
        
        t1 = slide.shapes.add_textbox(x, y + Inches(0.3), Inches(1.9), Inches(0.6))
        tf1 = t1.text_frame
        tf1.text = val
        p1 = tf1.paragraphs[0]
        p1.font.size = Pt(40)
        p1.font.bold = True
        p1.font.color.rgb = GOLDEN_YELLOW
        p1.alignment = PP_ALIGN.CENTER
        
        t2 = slide.shapes.add_textbox(x, y + Inches(1.0), Inches(1.9), Inches(0.4))
        tf2 = t2.text_frame
        tf2.text = lbl
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Rules & Tiers ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, NAVY_BLUE)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = "AUCTION OVERVIEW"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLDEN_YELLOW

    rules = [
        "10 College Teams • 11 Players Each",
        "Min 1 Wicket-Keeper • Min 3 Bowlers per team",
        "Exactly 4 Overseas players in Playing XI",
        "Recommended Purse: ₹1,20,00,000 per team"
    ]
    add_rules_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5), "AUCTION RULES", rules)

    tiers = [
        "10 ⭐⭐⭐⭐⭐ ₹2,00,00,000",
        "9 ⭐⭐⭐⭐½ ₹1,50,00,000",
        "8 ⭐⭐⭐⭐ ₹1,00,00,000",
        "7 ⭐⭐⭐½ ₹75,00,000",
        "6 ⭐⭐⭐ ₹50,00,000",
        "5 ⭐⭐½ ₹30,00,000",
        "4 ⭐⭐ ₹20,00,000"
    ]
    add_rules_box(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5), "BASE PRICE TIERS", tiers)

    # --- Slides 3+: Player Slides ---
    for player in players:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, NAVY_BLUE)

        # Bottom Bar Background
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))
        rect.fill.solid()
        rect.fill.fore_color.rgb = DARK_GREY_NAVY
        rect.line.width = 0

        # Player Name (Big)
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
        tf = name_box.text_frame
        tf.text = player['name'].upper()
        p = tf.paragraphs[0]
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Team/Role Divider Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(5), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLDEN_YELLOW

        # Team/Role Label
        role_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(10), Inches(0.5))
        tf = role_box.text_frame
        tf.text = f"{player['role']} | {player.get('overseas', 'Indian')} | {player.get('country', 'India')}"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = GOLDEN_YELLOW

        # Stats Area
        stats = stats_data.get(str(player['id']), {})
        # Formatting helper
        def get_stat(val, fallback="N/A"):
            return str(val) if val not in ["0", 0, "N/A", None] else fallback

        stats_display = [
            ("TOTAL RUNS", get_stat(stats.get('totalRuns'))),
            ("TOTAL WICKETS", get_stat(stats.get('totalWickets'))),
            ("STRIKE RATE", get_stat(stats.get('battingSR'))),
            ("ECONOMY", get_stat(stats.get('bowlingEconomy'))),
            ("BEST SCORE", get_stat(stats.get('bestScore'))),
            ("BEST BOWLING", get_stat(stats.get('bestBowling')))
        ]

        # Draw Stats Table (Gold Borders)
        for i, (label, val) in enumerate(stats_display):
            row = i // 2
            col = i % 2
            x = Inches(0.6 + (col * 3.2))
            y = Inches(3.0 + (row * 1.2))
            
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.0), Inches(1.0))
            box.line.color.rgb = GOLDEN_YELLOW
            box.fill.background()

            t1 = slide.shapes.add_textbox(x, y + Inches(0.1), Inches(3.0), Inches(0.3))
            tf1 = t1.text_frame
            tf1.text = label
            p1 = tf1.paragraphs[0]
            p1.font.size = Pt(12)
            p1.font.color.rgb = GOLDEN_YELLOW
            p1.alignment = PP_ALIGN.CENTER

            t2 = slide.shapes.add_textbox(x, y + Inches(0.4), Inches(3.0), Inches(0.5))
            tf2 = t2.text_frame
            tf2.text = val
            p2 = tf2.paragraphs[0]
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.CENTER

        # Footer Details (Base Price & Style)
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.5))
        tf = footer_box.text_frame
        tf.text = f"BASE PRICE: {player.get('basePrice', 'N/A')}  |  STYLE: {player.get('batting_hand', 'N/A')} & {player.get('bowling_style', 'N/A')}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Image Framing (Right side)
        img_name = player.get('image_file') or f"{player['id']}.png"
        img_path = os.path.join(images_dir, img_name)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), height=Inches(5.0))
            except: pass

    prs.save(output_path)
    print(f"Generated: {output_path}")

if __name__ == "__main__":
    generate_ppt()
