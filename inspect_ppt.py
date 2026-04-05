from pptx import Presentation

prs = Presentation('D:/ipl auction/IPL2026_College_Auction_Final.pptx')
slide = prs.slides[1]

with open('slide_shapes.txt', 'w', encoding='utf-8') as f:
    f.write(f"Total Shapes on Slide 1: {len(slide.shapes)}\n")
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            f.write(f"Shape {i} TEXT:\n{shape.text}\n---\n")
        else:
            f.write(f"Shape {i}: {shape.shape_type} / {shape.name}\n")
