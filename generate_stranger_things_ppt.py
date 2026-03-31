import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def extract_mock_stats(file_path):
    if not os.path.exists(file_path):
        return {}
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    match = re.search(r'export const mockPlayerStats = ({[\s\S]*?});', content)
    if not match: return {}
    js_obj = match.group(1)
    json_str = re.sub(r'(\s)(\w+):', r'\1"\2":', js_obj)
    json_str = re.sub(r',\s*}', '}', json_str)
    try:
        return json.loads(json_str)
    except:
        return {}

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def generate_ppt():
    players_json_path = 'client/public/players.json'
    mock_stats_path = 'client/src/mockStats.js'
    images_dir = 'client/public/players'
    output_path = 'Stranger_Things_Auction_Catalog.pptx'

    with open(players_json_path, 'r', encoding='utf-8') as f:
        players = json.load(f)
    stats_data = extract_mock_stats(mock_stats_path)

    prs = Presentation()
    # Slide dimensions: 16:9 
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Introduction ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    apply_background(slide, RGBColor(0, 0, 0))

    # Title "IPL AUCTION 2026"
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(13.333), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.text = "IPL AUCTION 2026"
    p = title_tf.paragraphs[0]
    p.font.name = 'Arial Black'
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0) # Neon Red
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(3.8), Inches(13.333), Inches(1))
    sub_tf = sub_box.text_frame
    sub_tf.text = "STRANGER THINGS AUCTION NIGHT\nEnter the Upside Down Auction"
    p = sub_tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 200, 255) # Ice Blue
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Rules ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, RGBColor(0, 0, 0))

    rule_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    rt_tf = rule_title.text_frame
    rt_tf.text = "THE RULES OF THE UPSIDE DOWN"
    p = rt_tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)

    rules = [
        "• Each team has a fixed budget (Don't let the Demogorgon take it)",
        "• Players will be auctioned one by one into the portal",
        "• Highest bidder claims the soul of the player",
        "• Breaking the budget results in catastrophic consequences",
        "• Bidding rounds are time-limited. Tick... Tock..."
    ]
    
    rules_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    rules_tf = rules_box.text_frame
    for r in rules:
        p = rules_tf.add_paragraph()
        p.text = r
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(15)

    # --- Slides 3+: Player Cards ---
    # We will loop through players. (Limiting to first 200 for safe file size)
    for player in players:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, RGBColor(5, 5, 15)) # Very dark blue/black

        # Glowing Header
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
        name_tf = name_box.text_frame
        name_tf.text = player['name'].upper()
        p = name_tf.paragraphs[0]
        p.font.name = 'Arial Black'
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)
        p.alignment = PP_ALIGN.LEFT

        # Team & Role Subheader
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(10), Inches(0.5))
        sub_tf = sub_box.text_frame
        sub_tf.text = f"{player['role']} | {player.get('team', 'FREE AGENT')} | {player.get('country', 'INDIAN')}"
        p = sub_tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 229, 255) # Cyan neon

        # Player Details (Base Price & Origin)
        base_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5), Inches(0.5))
        base_tf = base_box.text_frame
        base_tf.text = f"BASE PRICE: {player.get('basePrice', 'N/A')}"
        p = base_tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(200, 200, 200)

        # Stats Area - Using placeholders for the extraction rule
        stats = stats_data.get(str(player['id']), {})
        role = player.get('role', 'All-Rounder').lower()
        
        # Extract values
        runs = stats.get('totalRuns', '0')
        wkts = stats.get('totalWickets', '0')
        sr = stats.get('battingSR', 'N/A')
        econ = stats.get('bowlingEconomy', 'N/A')
        best_bat = stats.get('bestScore', '0')
        best_bowl = stats.get('bestBowling', '0/0')

        # Role-Based Extraction Logic for Slide Display
        display_runs = runs if runs not in ['0', 0] else "N/A"
        display_wkts = wkts if wkts not in ['0', 0] else "N/A"
        display_sr = sr if sr != "N/A" else "N/A"
        display_econ = econ if econ != "N/A" else "N/A"
        display_best_bat = best_bat if best_bat not in ['0', 0] else "N/A"
        display_best_bowl = best_bowl if best_bowl not in ['0/0', 'N/A'] else "N/A"

        # Apply Optional/Required Logic
        if "batsman" in role or "batter" in role or "wicket" in role:
            # Bowling is optional, show N/A if zero
            pass
        elif "bowler" in role:
            # Batting is optional, show N/A if zero
            pass
        
        stats_labels = [
            ("TOTAL RUNS", display_runs),
            ("TOTAL WICKETS", display_wkts),
            ("STRIKE RATE", display_sr),
            ("ECONOMY", display_econ),
            ("BEST BATTING", display_best_bat),
            ("BEST BOWLING", display_best_bowl)
        ]

        # Draw Stats Card (Table-like look)
        left_start = 5.0
        top_start = 2.5
        box_w = 3.5
        box_h = 1.2
        
        for i, (label, val) in enumerate(stats_labels):
            row = i // 2
            col = i % 2
            
            # Draw box border
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                        Inches(left_start + col * (box_w + 0.2)), 
                                        Inches(top_start + row * (box_h + 0.2)), 
                                        Inches(box_w), Inches(box_h))
            rect.fill.background()
            line = rect.line
            line.color.rgb = RGBColor(255, 0, 0) if i%2==0 else RGBColor(0, 229, 255)
            line.width = Pt(2)
            
            # Content
            txt = slide.shapes.add_textbox(Inches(left_start + col * (box_w + 0.2)), 
                                           Inches(top_start + row * (box_h + 0.2)), 
                                           Inches(box_w), Inches(box_h))
            txt_tf = txt.text_frame
            txt_tf.word_wrap = True
            
            p1 = txt_tf.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(14)
            p1.font.color.rgb = RGBColor(180, 180, 180)
            p1.alignment = PP_ALIGN.CENTER
            
            p2 = txt_tf.add_paragraph()
            p2.text = str(val) if val != "0" and val != 0 else "N/A"
            p2.font.size = Pt(32)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.alignment = PP_ALIGN.CENTER

        # Playing Style (Footer)
        style_footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        style_tf = style_footer.text_frame
        style_tf.text = f"PLAYING STYLE: {player.get('batting_hand', 'N/A')} & {player.get('bowling_style', 'N/A')}"
        p = style_tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(255, 0, 0)

        # Image (if exists)
        img_name = player.get('image_file') or f"{player['id']}.png"
        img_path = os.path.join(images_dir, img_name)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(2.3), height=Inches(4.0))
            except: pass

    prs.save(output_path)
    print(f"Successfully generated: {output_path}")

if __name__ == "__main__":
    generate_ppt()
